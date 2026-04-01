import asyncio
import base64
import json
import logging
import os
import time
import uuid
from datetime import datetime, timezone
from io import BytesIO
from typing import Optional
import zipfile

import anthropic
from fastapi import FastAPI, File, Form, HTTPException, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, StreamingResponse
from fastapi.staticfiles import StaticFiles
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

# In-memory job store: job_id -> {"image_data": bytes, "mime_type": str, ...}
jobs: dict[str, dict] = {}
# In-memory download store: token -> bytes
downloads: dict[str, bytes] = {}

# ── Logging setup ──────────────────────────────────────────────────────────────

logging.basicConfig(level=logging.INFO, format="%(asctime)s %(levelname)s %(message)s")
logger = logging.getLogger("snap_to_slide")

LOGS_DIR      = os.path.join(os.path.dirname(os.path.abspath(__file__)), "logs")
TEMPLATES_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "templates")

# ── Shape catalog ──────────────────────────────────────────────────────────────
#
# This is the single source of truth for every element type Claude can emit.
# It drives both the extraction prompt (automatically) and the renderer.
#
# To add a new shape: append an entry here — no other code changes needed.
#
# Fields:
#   desc         → sentence shown to Claude describing when to use this type
#   render       → rendering path ("text" | "stat" | "icon" | "connector")
#                  omit for auto-shapes (uses the "mso" integer instead)
#   mso          → OOXML / VBA MSO integer for auto-shapes (drives add_shape())
#   outline_only → True: transparent fill, colored border (uses fill_color as border)

SHAPE_CATALOG: dict[str, dict] = {
    # ── Text containers ───────────────────────────────────────────────────────
    "title":    {"render": "text",
                 "desc": "Primary slide heading — large, bold, usually at the top"},
    "subtitle": {"render": "text",
                 "desc": "Secondary heading below the title, slightly smaller and lighter"},
    "text_block": {"render": "text",
                   "desc": "Paragraph or bullet list. Separate lines with \\n; prefix bullets with '• '"},
    "stat":     {"render": "stat",
                 "desc": "Large numeric / KPI value with a smaller caption beneath it"},
    "icon":     {"render": "icon",
                 "desc": "Icon or symbol — describe it in the 'text' field; rendered as a centred label"},
    # ── Line connectors ───────────────────────────────────────────────────────
    "arrow":    {"render": "connector",
                 "desc": "Directed line connector between two points (x1,y1) → (x2,y2)"},
    # ── Filled geometric shapes ───────────────────────────────────────────────
    "shape_rectangle":      {"mso": 1,
                              "desc": "Filled rectangle or square"},
    "shape_rounded_rect":   {"mso": 5,
                              "desc": "Rectangle with rounded corners (filled)"},
    "shape_oval":           {"mso": 9,
                              "desc": "Ellipse or circle (filled)"},
    "shape_diamond":        {"mso": 4,
                              "desc": "Diamond / rhombus (filled)"},
    "shape_triangle":       {"mso": 7,
                              "desc": "Isosceles triangle (filled)"},
    "shape_parallelogram":  {"mso": 2,
                              "desc": "Parallelogram / slanted rectangle (filled)"},
    "shape_trapezoid":      {"mso": 3,
                              "desc": "Trapezoid (filled)"},
    # ── Outline-only shapes ───────────────────────────────────────────────────
    "shape_rectangle_outline":    {"mso": 1, "outline_only": True,
                                   "desc": "Rectangle with a colored border and no fill"},
    "shape_rounded_rect_outline": {"mso": 5, "outline_only": True,
                                   "desc": "Rounded rectangle with a colored border and no fill"},
    "shape_oval_outline":         {"mso": 9, "outline_only": True,
                                   "desc": "Ellipse / circle with a colored border and no fill"},
    # ── Process / flow shapes ─────────────────────────────────────────────────
    "shape_chevron":     {"mso": 52,
                          "desc": "Right-pointing chevron — use for mid-sequence process steps"},
    "shape_home_plate":  {"mso": 26,
                          "desc": "Pentagon / home-plate — first step in a left-to-right process flow"},
    "shape_arrow_right": {"mso": 13,
                          "desc": "Solid block right-pointing arrow (a filled shape, not a connector)"},
    "shape_arrow_left":  {"mso": 34,
                          "desc": "Solid block left-pointing arrow (filled shape)"},
    "shape_arrow_up":    {"mso": 35,
                          "desc": "Solid block upward-pointing arrow (filled shape)"},
    "shape_arrow_down":  {"mso": 36,
                          "desc": "Solid block downward-pointing arrow (filled shape)"},
}

# ── Slide theme palette catalog ────────────────────────────────────────────────
#
# To add a new theme: append an entry here.
# If a file  templates/<name>.pptx  exists, it is used as the base presentation
# (provides fonts and master layouts); otherwise a blank presentation is created
# and these palette colours are applied.

SLIDE_THEMES: dict[str, dict] = {
    "Default":    {"bg": "f5f0e8", "accent": "c84b31", "text": "2c2c2c",
                   "desc": "Warm off-white with terracotta accent"},
    "Atlas":      {"bg": "1e2d3d", "accent": "4a90d9", "text": "deeaf7",
                   "desc": "Dark navy with sky-blue accent"},
    "Celestial":  {"bg": "1a1a2e", "accent": "9b59b6", "text": "e8e8f0",
                   "desc": "Deep space blue with purple accent"},
    "Madison":    {"bg": "ffffff", "accent": "1f4e79", "text": "242424",
                   "desc": "Clean white with deep-navy accent"},
    "Retrospect": {"bg": "f2f2f2", "accent": "d04a35", "text": "404040",
                   "desc": "Light grey with bold red accent"},
    "Slate":      {"bg": "2d3748", "accent": "38b2ac", "text": "e2e8f0",
                   "desc": "Dark charcoal with teal highlight"},
    "Corporate":  {"bg": "ffffff", "accent": "0070c0", "text": "222222",
                   "desc": "Clean white with professional blue"},
    "Organic":    {"bg": "e8f5e9", "accent": "2e7d32", "text": "1b3a1e",
                   "desc": "Soft sage green — natural and calm"},
    "Solarized":  {"bg": "fdf6e3", "accent": "b58900", "text": "657b83",
                   "desc": "Warm parchment with gold accent"},
}


def _ts() -> str:
    """ISO-8601 timestamp with UTC timezone."""
    return datetime.now(timezone.utc).isoformat()


def _mime_to_ext(mime_type: str) -> str:
    return {
        "image/jpeg": ".jpg",
        "image/jpg":  ".jpg",
        "image/png":  ".png",
        "image/webp": ".webp",
        "image/gif":  ".gif",
        "image/bmp":  ".bmp",
    }.get((mime_type or "").lower(), ".img")


def _job_log_dir(job_id: str) -> str:
    path = os.path.join(LOGS_DIR, job_id)
    os.makedirs(path, exist_ok=True)
    return path


def _write_job_log(job_id: str, log: dict) -> None:
    """Serialise the log dict to logs/<job_id>/log.json."""
    try:
        path = os.path.join(_job_log_dir(job_id), "log.json")
        with open(path, "w", encoding="utf-8") as fh:
            json.dump(log, fh, indent=2, ensure_ascii=False, default=str)
        logger.info("job=%s log written to %s", job_id, path)
    except Exception as exc:
        logger.warning("job=%s failed to write log: %s", job_id, exc)


# ── Template / theme helpers ───────────────────────────────────────────────────

def _clear_template_slides(prs: Presentation) -> None:
    """Remove all content slides from a template, keeping master/layouts."""
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        try:
            prs.part.drop_rel(sldId.rId)
        except Exception:
            pass
        sldIdLst.remove(sldId)


def _get_base_presentation(theme_name: str) -> tuple[Presentation, dict]:
    """Return (prs, palette) for the named theme.

    If  templates/<theme_name>.pptx  exists it is used as the base (provides
    fonts and slide master); otherwise a blank Presentation is created.
    The palette dict supplies default bg / accent / text colours.
    """
    palette = SLIDE_THEMES.get(theme_name, SLIDE_THEMES["Default"])
    tpl_path = os.path.join(TEMPLATES_DIR, f"{theme_name}.pptx")
    if os.path.isfile(tpl_path):
        try:
            prs = Presentation(tpl_path)
            _clear_template_slides(prs)
            prs.slide_width  = Inches(10)
            prs.slide_height = Inches(5.625)
            logger.info("theme=%s loaded from %s", theme_name, tpl_path)
            return prs, palette
        except Exception as exc:
            logger.warning("theme=%s failed to load template: %s — using blank", theme_name, exc)

    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.625)
    return prs, palette


def _safe_add_shape(slide, mso_id: int, x: int, y: int, w: int, h: int):
    """add_shape() with a fallback to RECTANGLE if the MSO id is invalid."""
    try:
        return slide.shapes.add_shape(mso_id, x, y, w, h)
    except (ValueError, Exception):
        return slide.shapes.add_shape(1, x, y, w, h)  # rectangle fallback


_FALLBACK_RGB = RGBColor(0x80, 0x80, 0x80)

def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex string (with or without #) to RGBColor. Falls back to gray on any error."""
    try:
        h = str(hex_color).lstrip("#").strip()
        if len(h) != 6:
            return _FALLBACK_RGB
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    except Exception:
        return _FALLBACK_RGB


def add_top_bar(slide, width, accent_hex: str):
    """Add full-width accent bar at top (~5px tall)."""
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(0), width, Pt(6),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = hex_to_rgb(accent_hex)
    bar.line.fill.background()


def add_footer(slide, width, height, text_hex: str):
    """Add 'Snap to Slide' footer bottom-left."""
    tf_box = slide.shapes.add_textbox(
        Inches(0.2), height - Pt(20), Inches(3), Pt(20)
    )
    tf = tf_box.text_frame
    tf.text = "Snap to Slide"
    p = tf.paragraphs[0]
    run = p.runs[0]
    run.font.size = Pt(8)
    run.font.color.rgb = hex_to_rgb(text_hex)
    run.font.bold = False


def rect_border_point(cx, cy, hw, hh, dx, dy):
    """Return the point on a rectangle's border in direction (dx, dy) from center.
    cx, cy = center; hw, hh = half-width, half-height; dx, dy = direction."""
    if dx == 0 and dy == 0:
        return cx, cy
    ts = []
    if dx > 0:
        ts.append(hw / dx)
    elif dx < 0:
        ts.append(-hw / dx)
    if dy > 0:
        ts.append(hh / dy)
    elif dy < 0:
        ts.append(-hh / dy)
    t = min(ts)
    return cx + t * dx, cy + t * dy


def build_bullets_slide(prs, data: dict) -> None:
    theme = data["theme"]
    content = data["content"]
    title_text = data.get("title", "")

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(theme["bg"])

    add_top_bar(slide, slide_width, theme["accent"])

    # Left accent bar
    bar = slide.shapes.add_shape(
        1,
        Inches(0.3), Inches(0.5), Pt(5), Inches(4.5),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = hex_to_rgb(theme["accent"])
    bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.65), Inches(0.45), Inches(8.8), Inches(0.7)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = hex_to_rgb(theme["primary"])

    # Bullets
    bullets = content.get("bullets", [])
    bullet_box = slide.shapes.add_textbox(
        Inches(0.65), Inches(1.3), Inches(8.8), Inches(3.4)
    )
    tf = bullet_box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"• {bullet}"
        run.font.size = Pt(16)
        run.font.color.rgb = hex_to_rgb(theme["text"])

    # Note bar at bottom
    note = content.get("note", "")
    if note:
        note_box = slide.shapes.add_shape(
            1,
            Inches(0), slide_height - Inches(0.55), slide_width, Inches(0.45),
        )
        note_box.fill.solid()
        note_box.fill.fore_color.rgb = hex_to_rgb(theme["accent"])
        note_box.line.fill.background()
        tf2 = note_box.text_frame
        tf2.margin_left = Inches(0.2)
        p = tf2.paragraphs[0]
        run = p.add_run()
        run.text = note
        run.font.size = Pt(10)
        run.font.color.rgb = hex_to_rgb(theme["bg"])

    add_footer(slide, slide_width, slide_height, theme["text"])


def build_two_column_slide(prs, data: dict) -> None:
    theme = data["theme"]
    content = data["content"]
    title_text = data.get("title", "")

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(theme["bg"])

    add_top_bar(slide, slide_width, theme["accent"])

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.4), Inches(0.35), Inches(9.2), Inches(0.65)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = hex_to_rgb(theme["primary"])

    # Horizontal divider
    divider = slide.shapes.add_shape(
        1,
        Inches(0.4), Inches(1.1), Inches(9.2), Pt(2),
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = hex_to_rgb(theme["accent"])
    divider.line.fill.background()

    # Left column
    left_text = content.get("left_column", "")
    left_box = slide.shapes.add_textbox(
        Inches(0.4), Inches(1.25), Inches(4.4), Inches(3.8)
    )
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(left_text.split("\n") if left_text else []):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(14)
        run.font.color.rgb = hex_to_rgb(theme["text"])

    # Vertical separator
    vsep = slide.shapes.add_shape(
        1,
        Inches(4.95), Inches(1.15), Pt(2), Inches(3.9),
    )
    vsep.fill.solid()
    vsep.fill.fore_color.rgb = hex_to_rgb(theme["accent"])
    vsep.line.fill.background()

    # Right column
    right_text = content.get("right_column", "")
    right_box = slide.shapes.add_textbox(
        Inches(5.15), Inches(1.25), Inches(4.4), Inches(3.8)
    )
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(right_text.split("\n") if right_text else []):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(14)
        run.font.color.rgb = hex_to_rgb(theme["text"])

    add_footer(slide, slide_width, slide_height, theme["text"])


def build_key_stats_slide(prs, data: dict) -> None:
    theme = data["theme"]
    content = data["content"]
    title_text = data.get("title", "")

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(theme["bg"])

    add_top_bar(slide, slide_width, theme["accent"])

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.4), Inches(0.35), Inches(9.2), Inches(0.65)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = hex_to_rgb(theme["primary"])

    stats = content.get("stats", [])
    if not stats:
        add_footer(slide, slide_width, slide_height, theme["text"])
        return

    n = len(stats)
    card_w = Inches(min(2.5, 9.0 / n))
    card_h = Inches(2.0)
    total_w = card_w * n
    start_x = (slide_width - total_w) / 2
    card_y = (slide_height - card_h) / 2 + Inches(0.2)

    for i, stat in enumerate(stats):
        cx = start_x + i * card_w

        # Card background
        card = slide.shapes.add_shape(
            1,
            cx + Inches(0.05), card_y, card_w - Inches(0.1), card_h,
        )
        card.fill.solid()
        card.fill.fore_color.rgb = hex_to_rgb(theme["primary"])
        card.line.fill.background()

        # Stat value
        val_box = slide.shapes.add_textbox(
            cx + Inches(0.05), card_y + Inches(0.15),
            card_w - Inches(0.1), Inches(1.1)
        )
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = stat.get("value", "")
        run.font.size = Pt(36)
        run.font.bold = True
        run.font.color.rgb = hex_to_rgb(theme["bg"])

        # Stat label
        lbl_box = slide.shapes.add_textbox(
            cx + Inches(0.05), card_y + Inches(1.3),
            card_w - Inches(0.1), Inches(0.6)
        )
        tf = lbl_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = stat.get("label", "")
        run.font.size = Pt(12)
        run.font.color.rgb = hex_to_rgb(theme["bg"])

    add_footer(slide, slide_width, slide_height, theme["text"])


def build_title_content_slide(prs, data: dict) -> None:
    theme = data["theme"]
    content = data["content"]
    title_text = data.get("title", "")

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(theme["bg"])

    add_top_bar(slide, slide_width, theme["accent"])

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.35), Inches(9.0), Inches(0.7)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(30)
    run.font.bold = True
    run.font.color.rgb = hex_to_rgb(theme["primary"])

    # Subtitle
    subtitle = content.get("subtitle", "")
    if subtitle:
        sub_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.1), Inches(9.0), Inches(0.45)
        )
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = subtitle
        run.font.size = Pt(16)
        run.font.italic = True
        run.font.color.rgb = hex_to_rgb(theme["accent"])

    # Main text panel
    main_text = content.get("main_text", "")
    panel_y = Inches(1.65) if subtitle else Inches(1.2)
    panel_h = Inches(3.0) if content.get("note") else Inches(3.4)

    panel = slide.shapes.add_shape(
        1,
        Inches(0.5), panel_y, Inches(9.0), panel_h,
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = hex_to_rgb(theme["primary"])
    panel.line.fill.background()
    tf2 = panel.text_frame
    tf2.word_wrap = True
    tf2.margin_left = Inches(0.2)
    tf2.margin_top = Inches(0.15)
    tf2.margin_right = Inches(0.2)
    p = tf2.paragraphs[0]
    run = p.add_run()
    run.text = main_text
    run.font.size = Pt(15)
    run.font.color.rgb = hex_to_rgb(theme["bg"])

    # Note bar
    note = content.get("note", "")
    if note:
        note_box = slide.shapes.add_shape(
            1,
            Inches(0), slide_height - Inches(0.55), slide_width, Inches(0.45),
        )
        note_box.fill.solid()
        note_box.fill.fore_color.rgb = hex_to_rgb(theme["accent"])
        note_box.line.fill.background()
        tf3 = note_box.text_frame
        tf3.margin_left = Inches(0.2)
        p = tf3.paragraphs[0]
        run = p.add_run()
        run.text = note
        run.font.size = Pt(10)
        run.font.color.rgb = hex_to_rgb(theme["bg"])

    add_footer(slide, slide_width, slide_height, theme["text"])


def build_diagram_slide(prs, data: dict) -> None:
    """Render a geometry-aware diagram slide: boxes (nodes) connected by arrows (edges)."""
    theme = data["theme"]
    content = data["content"]
    title_text = data.get("title", "")

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Content area bounds (nodes' normalized coords map into this region)
    AREA_X = Inches(0.3)
    AREA_Y = Inches(0.9)
    AREA_W = Inches(9.4)
    AREA_H = Inches(4.2)

    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(theme["bg"])

    add_top_bar(slide, slide_width, theme["accent"])

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.4), Inches(0.2), Inches(9.2), Inches(0.6)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = hex_to_rgb(theme["primary"])

    nodes = content.get("nodes", [])
    edges = content.get("edges", [])

    # Compute absolute positions and build center map for edge drawing
    node_map: dict[str, dict] = {}
    for node in nodes:
        ax = AREA_X + node["x"] * AREA_W
        ay = AREA_Y + node["y"] * AREA_H
        aw = node["w"] * AREA_W
        ah = node["h"] * AREA_H
        node_map[node["id"]] = {
            "cx": ax + aw / 2,
            "cy": ay + ah / 2,
            "hw": aw / 2,
            "hh": ah / 2,
            "x": ax, "y": ay, "w": aw, "h": ah,
        }

    # Draw edges first (behind nodes)
    for edge in edges:
        fid = edge.get("from")
        tid = edge.get("to")
        if fid not in node_map or tid not in node_map:
            continue
        fn = node_map[fid]
        tn = node_map[tid]
        dx = tn["cx"] - fn["cx"]
        dy = tn["cy"] - fn["cy"]
        # Exit point from source node border
        x1, y1 = rect_border_point(fn["cx"], fn["cy"], fn["hw"], fn["hh"], dx, dy)
        # Entry point at target node border
        x2, y2 = rect_border_point(tn["cx"], tn["cy"], tn["hw"], tn["hh"], -dx, -dy)

        try:
            connector = slide.shapes.add_connector(
                MSO_CONNECTOR_TYPE.STRAIGHT, x1, y1, x2, y2
            )
            connector.line.color.rgb = hex_to_rgb(theme["accent"])
            connector.line.width = Pt(1.5)
            # Add arrowhead at the end using OxmlElement (correct namespace handling)
            ln = connector.line._ln
            if ln is not None:
                tail_end = OxmlElement("a:tailEnd")
                tail_end.set("type", "arrow")
                tail_end.set("w", "med")
                tail_end.set("len", "med")
                ln.append(tail_end)
        except Exception:
            pass  # Connector drawing is best-effort

        # Edge label (if any)
        label = edge.get("label", "")
        if label:
            lx = (x1 + x2) / 2 - Inches(0.4)
            ly = (y1 + y2) / 2 - Pt(8)
            lbl = slide.shapes.add_textbox(lx, ly, Inches(0.8), Pt(16))
            tf = lbl.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = label
            run.font.size = Pt(9)
            run.font.italic = True
            run.font.color.rgb = hex_to_rgb(theme["accent"])

    # Draw nodes (on top of edges)
    for node in nodes:
        nm = node_map[node["id"]]
        shape = slide.shapes.add_shape(
            1,  # MSO_AUTO_SHAPE_TYPE.RECTANGLE (same as original code uses)
            nm["x"], nm["y"], nm["w"], nm["h"],
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(theme["primary"])
        shape.line.color.rgb = hex_to_rgb(theme["accent"])
        shape.line.width = Pt(1)

        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.05)
        tf.margin_right = Inches(0.05)
        tf.margin_top = Inches(0.03)
        tf.margin_bottom = Inches(0.03)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = node.get("text", "")
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = hex_to_rgb(theme["bg"])

    add_footer(slide, slide_width, slide_height, theme["text"])


def build_table_slide(prs, data: dict) -> None:
    """Render a structured table slide."""
    theme = data["theme"]
    content = data["content"]
    title_text = data.get("title", "")

    headers = content.get("headers", [])
    rows = content.get("rows", [])

    if not headers and not rows:
        build_bullets_slide(prs, data)
        return

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(theme["bg"])

    add_top_bar(slide, slide_width, theme["accent"])

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.4), Inches(0.2), Inches(9.2), Inches(0.6)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = hex_to_rgb(theme["primary"])

    num_cols = len(headers) if headers else (len(rows[0]) if rows else 1)
    num_rows = len(rows) + (1 if headers else 0)

    table_shape = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.4), Inches(0.9),
        Inches(9.2), Inches(4.3),
    )
    table = table_shape.table

    row_offset = 0
    if headers:
        for j, header in enumerate(headers[:num_cols]):
            cell = table.cell(0, j)
            cell.text = str(header)
            cell.fill.solid()
            cell.fill.fore_color.rgb = hex_to_rgb(theme["primary"])
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.runs[0] if p.runs else p.add_run()
            run.font.bold = True
            run.font.size = Pt(12)
            run.font.color.rgb = hex_to_rgb(theme["bg"])
        row_offset = 1

    for i, row in enumerate(rows):
        for j, cell_text in enumerate(list(row)[:num_cols]):
            cell = table.cell(i + row_offset, j)
            cell.text = str(cell_text)
            cell.fill.solid()
            cell.fill.fore_color.rgb = hex_to_rgb(theme["bg"])
            p = cell.text_frame.paragraphs[0]
            run = p.runs[0] if p.runs else p.add_run()
            run.font.size = Pt(11)
            run.font.color.rgb = hex_to_rgb(theme["text"])

    add_footer(slide, slide_width, slide_height, theme["text"])


# ── Layered slide builder ──────────────────────────────────────────────────────

def _safe_float(val, default: float = 0.0) -> float:
    try:
        return float(val)
    except Exception:
        return default


def _render_shape_element(slide, el: dict, sw: int, sh: int) -> None:
    """Render any auto-shape from SHAPE_CATALOG (those with an 'mso' key)."""
    etype = el.get("type", "shape_rectangle")
    entry = SHAPE_CATALOG.get(etype, {})

    # Legacy support: old "shape" type with a shape_type sub-field
    if etype == "shape":
        mso = {"rect": 1, "rounded_rect": 5, "oval": 9}.get(el.get("shape_type", "rect"), 1)
        outline_only = False
    else:
        mso = entry.get("mso", 1)
        outline_only = entry.get("outline_only", False)

    x = int(_safe_float(el.get("x", 0)) * sw)
    y = int(_safe_float(el.get("y", 0)) * sh)
    w = max(int(_safe_float(el.get("w", 0.1)) * sw), int(Pt(10)))
    h = max(int(_safe_float(el.get("h", 0.05)) * sh), int(Pt(10)))

    shape = _safe_add_shape(slide, mso, x, y, w, h)

    fill_hex   = el.get("fill_color", "")
    border_hex = el.get("border_color", "")

    if outline_only:
        shape.fill.background()
        color_hex = fill_hex or border_hex
        if color_hex:
            shape.line.color.rgb = hex_to_rgb(color_hex)
            shape.line.width = Pt(1.5)
    else:
        if fill_hex:
            shape.fill.solid()
            shape.fill.fore_color.rgb = hex_to_rgb(fill_hex)
        else:
            shape.fill.background()
        if border_hex:
            shape.line.color.rgb = hex_to_rgb(border_hex)
            shape.line.width = Pt(1.5)
        else:
            shape.line.fill.background()

    text = el.get("text", "")
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left  = Inches(0.06)
        tf.margin_right = Inches(0.06)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.size  = Pt(_safe_float(el.get("font_size", 12), 12))
        run.font.bold  = bool(el.get("bold", True))
        run.font.color.rgb = hex_to_rgb(el.get("text_color", "ffffff"))


def _render_text_element(slide, el: dict, sw: int, sh: int, etype: str) -> None:
    """Render title, subtitle, or text_block as a text box."""
    x = int(_safe_float(el.get("x", 0)) * sw)
    y = int(_safe_float(el.get("y", 0)) * sh)
    w = max(int(_safe_float(el.get("w", 0.5)) * sw), int(Pt(10)))
    h = max(int(_safe_float(el.get("h", 0.1)) * sh), int(Pt(10)))

    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True

    align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
    alignment = align_map.get(el.get("align", "left"), PP_ALIGN.LEFT)
    color_rgb = hex_to_rgb(el.get("color", "333333"))

    default_size = {"title": 28, "subtitle": 20}.get(etype, 14)
    font_size = Pt(_safe_float(el.get("font_size", default_size), default_size))
    bold   = bool(el.get("bold",   etype == "title"))
    italic = bool(el.get("italic", etype == "subtitle"))

    lines = str(el.get("text", "")).split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = alignment
        run = p.add_run()
        run.text        = line
        run.font.size   = font_size
        run.font.bold   = bold
        run.font.italic = italic
        run.font.color.rgb = color_rgb


def _render_icon_element(slide, el: dict, sw: int, sh: int) -> None:
    """Render an icon placeholder — displayed as a centred text label."""
    x = int(_safe_float(el.get("x", 0)) * sw)
    y = int(_safe_float(el.get("y", 0)) * sh)
    w = max(int(_safe_float(el.get("w", 0.08)) * sw), int(Pt(20)))
    h = max(int(_safe_float(el.get("h", 0.08)) * sh), int(Pt(20)))

    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(el.get("text", ""))
    run.font.size  = Pt(_safe_float(el.get("font_size", 20), 20))
    run.font.color.rgb = hex_to_rgb(el.get("color", "555555"))


def _render_arrow_element(slide, el: dict, sw: int, sh: int, theme: dict) -> None:
    x1 = int(_safe_float(el.get("x1", 0)) * sw)
    y1 = int(_safe_float(el.get("y1", 0)) * sh)
    x2 = int(_safe_float(el.get("x2", 0.5)) * sw)
    y2 = int(_safe_float(el.get("y2", 0)) * sh)
    color_hex = el.get("color", "") or theme.get("accent", "c84b31")

    connector = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, x1, y1, x2, y2)
    connector.line.color.rgb = hex_to_rgb(color_hex)
    connector.line.width = Pt(1.5)
    ln = connector.line._ln
    if ln is not None:
        tail_end = OxmlElement("a:tailEnd")
        tail_end.set("type", "arrow")
        tail_end.set("w", "med")
        tail_end.set("len", "med")
        ln.append(tail_end)

    label = el.get("label", "")
    if label:
        lx = (x1 + x2) // 2 - int(Inches(0.4))
        ly = (y1 + y2) // 2 - int(Pt(8))
        lbl = slide.shapes.add_textbox(lx, ly, int(Inches(0.8)), int(Pt(16)))
        p = lbl.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        run.font.size = Pt(9)
        run.font.italic = True
        run.font.color.rgb = hex_to_rgb(color_hex)


def _render_stat_element(slide, el: dict, sw: int, sh: int) -> None:
    x = int(_safe_float(el.get("x", 0)) * sw)
    y = int(_safe_float(el.get("y", 0)) * sh)
    w = max(int(_safe_float(el.get("w", 0.2)) * sw), int(Pt(20)))
    h = max(int(_safe_float(el.get("h", 0.3)) * sh), int(Pt(20)))

    val_h = int(h * 0.6)
    lbl_h = max(h - val_h, int(Pt(10)))

    val_box = slide.shapes.add_textbox(x, y, w, val_h)
    p = val_box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(el.get("value", ""))
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = hex_to_rgb(el.get("value_color", "2c3e50"))

    lbl_box = slide.shapes.add_textbox(x, y + val_h, w, lbl_h)
    p2 = lbl_box.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = str(el.get("label", ""))
    run2.font.size = Pt(12)
    run2.font.color.rgb = hex_to_rgb(el.get("label_color", "7a7065"))


def build_layered_slide(prs, data: dict) -> None:
    """Build a slide from a flat list of positioned elements in back-to-front order.

    Dispatches each element to its renderer based on SHAPE_CATALOG. Adding a new
    shape type to the catalog is the only change needed to support it here.
    """
    raw_theme = data.get("theme", {})
    theme     = raw_theme if isinstance(raw_theme, dict) else {}
    elements  = data.get("elements", [])

    SLIDE_W = prs.slide_width
    SLIDE_H = prs.slide_height

    # Use blank layout (index 6); templates provide the master visuals
    layout_idx = min(6, len(prs.slide_layouts) - 1)
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = hex_to_rgb(theme.get("bg", "f5f0e8"))

    add_top_bar(slide, SLIDE_W, theme.get("accent", "c84b31"))

    for el in elements:
        etype = el.get("type", "")
        entry = SHAPE_CATALOG.get(etype, {})
        render_path = entry.get("render", "auto_shape" if "mso" in entry else None)

        try:
            if render_path == "text" or etype in ("title", "subtitle", "text_block"):
                _render_text_element(slide, el, SLIDE_W, SLIDE_H, etype)
            elif render_path == "connector" or etype == "arrow":
                _render_arrow_element(slide, el, SLIDE_W, SLIDE_H, theme)
            elif render_path == "stat" or etype == "stat":
                _render_stat_element(slide, el, SLIDE_W, SLIDE_H)
            elif render_path == "icon" or etype == "icon":
                _render_icon_element(slide, el, SLIDE_W, SLIDE_H)
            elif "mso" in entry or etype == "shape":
                # All auto-shapes from the catalog + legacy "shape" type
                _render_shape_element(slide, el, SLIDE_W, SLIDE_H)
            # Unknown types are silently skipped
        except Exception:
            pass  # best-effort: skip bad elements, don't abort the slide

    add_footer(slide, SLIDE_W, SLIDE_H, theme.get("text", "333333"))


def validate_pptx(pptx_bytes: bytes) -> tuple[bytes, bool]:
    """Validate PPTX bytes by reopening. Returns (bytes, used_fallback)."""
    buf = BytesIO(pptx_bytes)
    if not zipfile.is_zipfile(buf):
        return _make_fallback_pptx("Slide validation failed (invalid archive)"), True
    try:
        buf.seek(0)
        Presentation(buf)
        return pptx_bytes, False
    except Exception as e:
        return _make_fallback_pptx(f"Slide validation failed: {e}"), True


def _make_fallback_pptx(reason: str = "") -> bytes:
    """Return a minimal valid one-slide PPTX (used when the primary builder produces corrupt output)."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9.0), Inches(1.0))
    p = tb.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Slide could not be rendered"
    run.font.size = Pt(26)
    run.font.bold = True

    if reason:
        tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9.0), Inches(2.0))
        p2 = tb2.text_frame.paragraphs[0]
        run2 = p2.add_run()
        run2.text = reason
        run2.font.size = Pt(14)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def build_pptx(data: dict, theme_name: str = "Default") -> tuple[bytes, bool]:
    """Build PPTX on top of the requested theme/template. Returns (bytes, used_fallback)."""
    prs, palette = _get_base_presentation(theme_name)

    # Fill in any missing theme colours from the palette (Claude's colours take precedence)
    if not isinstance(data.get("theme"), dict):
        data["theme"] = {}
    for key in ("bg", "accent", "text"):
        data["theme"].setdefault(key, palette[key])

    layout = data.get("layout", "layered")
    builders = {
        "bullets": build_bullets_slide,
        "two_column": build_two_column_slide,
        "key_stats": build_key_stats_slide,
        "title_content": build_title_content_slide,
        "diagram": build_diagram_slide,
        "table": build_table_slide,
        "layered": build_layered_slide,
    }
    builder = builders.get(layout, build_layered_slide)
    builder(prs, data)

    buf = BytesIO()
    prs.save(buf)
    return validate_pptx(buf.getvalue())


# ── Prompts ────────────────────────────────────────────────────────────────────

STRUCTURE_PROMPT = """Analyze this image and identify its primary visual structure.
Return ONLY valid JSON (no markdown, no explanation):
{
  "structure_type": "diagram|table|chart|text",
  "notes": "one sentence describing the spatial layout and key visual features"
}

Structure types:
- "diagram": flowcharts, process flows, org charts, mind maps, network diagrams — shapes/boxes connected by arrows or lines
- "table": grids, matrices, comparison tables, spreadsheet-like rows and columns of data
- "chart": bar charts, pie charts, line graphs, scatter plots — numerical/statistical visualizations
- "text": written content, bullet lists, slides with paragraphs, notes, outlines, or any primarily textual content"""


DIAGRAM_EXTRACTION_PROMPT = """You are analyzing a diagram. Extract its structure as positioned nodes and directed edges, preserving the spatial layout exactly as it appears.

Return ONLY valid JSON:
{
  "title": "descriptive title of the diagram",
  "layout": "diagram",
  "theme": {
    "bg": "f5f0e8",
    "primary": "2c3e50",
    "accent": "c84b31",
    "text": "333333"
  },
  "content": {
    "nodes": [
      {"id": "n1", "text": "Node label", "x": 0.05, "y": 0.3, "w": 0.2, "h": 0.15}
    ],
    "edges": [
      {"from": "n1", "to": "n2", "label": ""}
    ]
  }
}

SPATIAL RULES — these are critical:
- x, y, w, h are normalized 0.0–1.0 relative to the diagram area
- x, y = top-left corner of the node; w, h = width and height of the node
- Mirror the exact spatial layout: if A appears left of B in the image, A must have a smaller x value
- Nodes at the top of the image get small y values; nodes at the bottom get larger y values
- Do NOT cluster all nodes at the same position — spread them to match the original
- Typical node: w=0.15–0.25, h=0.10–0.18; do not make nodes too small (min w=0.10, h=0.08)
- Nodes must not overlap: ensure rectangles [x, x+w] × [y, y+h] do not intersect
- Include ALL visible nodes and ALL arrows/connections
- edges: directed from source to destination; include label text if present on the arrow

Theme: choose colors that complement the content. All hex values 6 chars, NO # symbol."""


TABLE_EXTRACTION_PROMPT = """You are analyzing a table or grid. Extract it precisely as structured data.

Return ONLY valid JSON:
{
  "title": "table title or topic",
  "layout": "table",
  "theme": {
    "bg": "f8f9fa",
    "primary": "1a3a5c",
    "accent": "2980b9",
    "text": "2c3e50"
  },
  "content": {
    "headers": ["Column 1", "Column 2", "Column 3"],
    "rows": [
      ["row1col1", "row1col2", "row1col3"],
      ["row2col1", "row2col2", "row2col3"]
    ]
  }
}

Rules:
- Extract ALL visible rows and columns faithfully — do not summarize or skip rows
- headers: column header texts (empty array [] if there are no headers)
- rows: array of arrays; each inner array is one data row in order
- Preserve exact cell text content; do not paraphrase
- If a cell spans multiple columns, repeat the value across those columns
Theme: all hex values 6 chars, NO # symbol."""


CHART_EXTRACTION_PROMPT = """You are analyzing a chart or graph. Extract the key data points and metrics.

Return ONLY valid JSON:
{
  "title": "chart title",
  "layout": "key_stats",
  "theme": {
    "bg": "f5f0e8",
    "primary": "2c3e50",
    "accent": "c84b31",
    "text": "333333"
  },
  "content": {
    "stats": [
      {"value": "42%", "label": "Metric Name"}
    ],
    "note": "Chart type and key insight in one sentence"
  }
}

Rules:
- Extract the most important 2–5 data points as stats (the standout numbers, peaks, or totals)
- "note": describe the chart type (bar, pie, line…) and the primary trend or finding
- Keep value strings concise (e.g. "42%", "$1.2M", "3.7x")
Theme: all hex values 6 chars, NO # symbol."""


TEXT_EXTRACTION_PROMPT = """You are a PowerPoint slide designer. Analyze the image and extract its content, then design a single professional slide.

Return ONLY valid JSON (no markdown, no explanation) with this exact structure:
{
  "title": "slide title",
  "layout": "bullets",
  "theme": {
    "bg": "f5f0e8",
    "primary": "2c3e50",
    "accent": "c84b31",
    "text": "333333"
  },
  "content": {
    "main_text": "",
    "subtitle": "",
    "bullets": [],
    "left_column": "",
    "right_column": "",
    "stats": [],
    "note": ""
  }
}

Layout rules:
- "bullets": use for notes, lists, outlines — populate "bullets" array and optionally "note"
- "two_column": use for comparisons or two topics — populate "left_column" and "right_column"
- "key_stats": use for numbers, metrics, data — populate "stats" array as [{"value": "42%", "label": "Growth"}, ...]
- "title_content": use for single topic with explanation — populate "title", "subtitle", "main_text", optionally "note"

Theme: choose colors that suit the content. All hex values are 6 chars, NO # symbol.
Extract all visible text. Be faithful to the source content."""

EXTRACTION_PROMPTS = {
    "diagram": DIAGRAM_EXTRACTION_PROMPT,
    "table": TABLE_EXTRACTION_PROMPT,
    "chart": CHART_EXTRACTION_PROMPT,
    "text": TEXT_EXTRACTION_PROMPT,
}

def _build_extraction_prompt() -> str:
    """Build the Claude extraction prompt dynamically from SHAPE_CATALOG.

    Extending the catalog automatically extends the prompt Claude receives —
    no manual prompt editing required when adding a new shape type.
    """
    # ── Per-type field specs ──────────────────────────────────────────────────
    field_specs: dict[str, str] = {
        "title":      "x, y, w, h, font_size (int, default 28), bold (bool), "
                      "italic (bool), color (hex), align (\"left\"|\"center\"|\"right\")",
        "subtitle":   "x, y, w, h, font_size (int, default 20), bold (bool), "
                      "italic (bool, default true), color (hex), align",
        "text_block": "x, y, w, h, font_size (int, default 14), align, color (hex)",
        "stat":       "value (str), label (str), x, y, w, h, value_color (hex), label_color (hex)",
        "icon":       "text (symbol or short description), x, y, w, h, font_size (int), color (hex)",
        "arrow":      "x1, y1, x2, y2 (all normalized 0–1), color (hex), label (str or \"\")",
    }
    # Default field spec for auto-shapes (those with an "mso" key)
    auto_shape_fields = ("x, y, w, h, fill_color (hex or \"\"), border_color (hex or \"\"), "
                         "text (str or \"\"), text_color (hex), font_size (int), bold (bool)")

    # ── Build ELEMENT TYPES section from catalog ───────────────────────────────
    lines: list[str] = ["ELEMENT TYPES (use exactly these strings for the \"type\" field):\n"]
    for name, entry in SHAPE_CATALOG.items():
        desc   = entry.get("desc", "")
        fields = field_specs.get(name, auto_shape_fields)
        lines.append(f"{name} — {desc}:\n  {fields}\n")

    element_types_block = "\n".join(lines)

    return f"""You are a PowerPoint slide designer. Analyze this image and convert every \
visible element into a layered slide description. Identify ALL elements as separate \
positioned layers ordered back-to-front (backgrounds first, foreground labels last).

Return ONLY valid JSON (no markdown, no explanation):
{{
  "title": "descriptive slide title",
  "theme": {{
    "bg": "f5f0e8",
    "accent": "c84b31",
    "text": "333333"
  }},
  "elements": [
    {{"type": "shape_rectangle", "x": 0.0, "y": 0.0, "w": 1.0, "h": 0.05, "fill_color": "c84b31", "border_color": "", "text": "", "text_color": "ffffff", "font_size": 12, "bold": true}},
    {{"type": "title", "text": "Slide Heading", "x": 0.04, "y": 0.07, "w": 0.92, "h": 0.14, "font_size": 28, "bold": true, "italic": false, "color": "2c3e50", "align": "left"}},
    {{"type": "subtitle", "text": "Supporting subtitle", "x": 0.04, "y": 0.20, "w": 0.70, "h": 0.09, "font_size": 18, "bold": false, "italic": true, "color": "7a7065", "align": "left"}},
    {{"type": "text_block", "text": "• Bullet one\\n• Bullet two\\n• Bullet three", "x": 0.06, "y": 0.30, "w": 0.88, "h": 0.45, "font_size": 14, "align": "left", "color": "333333"}},
    {{"type": "shape_chevron", "x": 0.05, "y": 0.55, "w": 0.18, "h": 0.12, "fill_color": "2c3e50", "border_color": "", "text": "Step 1", "text_color": "ffffff", "font_size": 12, "bold": true}},
    {{"type": "arrow", "x1": 0.33, "y1": 0.61, "x2": 0.50, "y2": 0.61, "color": "c84b31", "label": ""}},
    {{"type": "stat", "value": "42%", "label": "Growth rate", "x": 0.10, "y": 0.70, "w": 0.22, "h": 0.20, "value_color": "2c3e50", "label_color": "7a7065"}}
  ]
}}

{element_types_block}
SHAPE SELECTION GUIDE:
- Process flows (A → B → C): use shape_home_plate for step 1, shape_chevron for remaining steps
- Boxes with borders only (no fill): use shape_rectangle_outline / shape_rounded_rect_outline
- Block arrows as shapes (not connectors): use shape_arrow_right / shape_arrow_left / shape_arrow_up / shape_arrow_down
- Directional connector lines with arrowheads: use arrow (x1,y1 → x2,y2)
- Circle / ellipse shapes: use shape_oval (filled) or shape_oval_outline (border only)

POSITIONING RULES:
- x, y, w, h are normalized 0.0–1.0 relative to the full slide (16:9, 10 × 5.625 in)
- x=0 is left edge, y=0 is top; x+w ≤ 1.0 and y+h ≤ 1.0
- Mirror the exact spatial layout of the original image (left/right and top/bottom positions)
- Include ALL visible text — do not omit or summarize content
- Minimum sizes: shapes w ≥ 0.08, h ≥ 0.06 · text boxes w ≥ 0.10, h ≥ 0.05
- For flowcharts: each box → appropriate shape element, each line with arrow → arrow element
- For charts / graphs: key numbers → stat elements; axis labels → text_block elements

THEME:
- bg: slide background hex colour (light for readability unless the original is dark)
- accent: colour for arrows, borders, and the accent bar at the top of the slide
- text: default body text colour
- All hex values: exactly 6 characters, NO # symbol. Use "" (empty string) for no fill/border.

OUTPUT: Return the JSON only — no explanation, no markdown code fences."""


LAYERED_EXTRACTION_PROMPT = _build_extraction_prompt()


# ── Claude Vision call ─────────────────────────────────────────────────────────

def _call_vision_model(
    image_bytes: bytes,
    mime_type: str,
    prompt: str,
    hint: str,
    api_key: str,
) -> str:
    """Call Claude Vision and return the raw response text."""
    b64 = base64.standard_b64encode(image_bytes).decode()
    client = anthropic.Anthropic(api_key=api_key)
    content = [
        {"type": "image",  "source": {"type": "base64", "media_type": mime_type, "data": b64}},
        {"type": "text",   "text": prompt},
    ]
    if hint:
        content.append({"type": "text", "text": f"User context: {hint}"})

    response = client.messages.create(
        model="claude-sonnet-4-5",
        max_tokens=4096,
        messages=[{"role": "user", "content": content}],
    )
    return response.content[0].text.strip()


def _build_preview(slide_data: dict) -> dict:
    """Extract human-readable text content from slide_data for UI display."""
    title = slide_data.get("title", "")
    items = []
    seen: set[str] = set()

    for el in slide_data.get("elements", []):
        etype = el.get("type", "")
        entry = SHAPE_CATALOG.get(etype, {})
        text = ""
        kind = "shape"

        if etype in ("title", "subtitle", "text_block"):
            text = str(el.get("text", "")).strip()
            kind = etype
        elif etype == "stat":
            value = str(el.get("value", "")).strip()
            label = str(el.get("label", "")).strip()
            text = f"{value}  {label}".strip() if (value or label) else ""
            kind = "stat"
        elif entry.get("render") == "icon" or etype == "icon":
            text = str(el.get("text", "")).strip()
            kind = "shape"
        elif "mso" in entry or etype == "shape":
            # Auto-shapes: only show if they have inner text
            text = str(el.get("text", "")).strip()
            kind = "shape"
        else:
            continue

        if text and text not in seen:
            seen.add(text)
            items.append({"kind": kind, "text": text})

    return {"title": title, "items": items}


async def process_job(job_id: str):
    """Generator that streams SSE events for a job and writes a detailed log to disk."""
    job = jobs.get(job_id)
    if not job:
        yield f"data: {json.dumps({'error': 'Job not found'})}\n\n"
        return

    image_data = job["image_data"]
    mime_type  = job["mime_type"]
    api_key    = job["api_key"]
    hint_text  = job.get("hint_text", "")
    theme_name = job.get("theme", "Default")

    def sse(payload: dict) -> str:
        return f"data: {json.dumps(payload)}\n\n"

    # ── Initialise log ─────────────────────────────────────────────────────────
    t_start = time.monotonic()
    log: dict = {
        "job_id":     job_id,
        "started_at": _ts(),
        "image_file": os.path.join(LOGS_DIR, job_id, f"image{_mime_to_ext(mime_type)}"),
        "mime_type":  mime_type,
        "hint_text":  hint_text,
        "steps":      {},
        "error":      None,
        "completed_at":      None,
        "total_duration_ms": None,
    }
    log["theme"] = theme_name
    logger.info("job=%s started mime=%s theme=%s hint=%r", job_id, mime_type, theme_name, hint_text)

    try:
        loop = asyncio.get_event_loop()

        # ── Step 1: Vision model call ──────────────────────────────────────────
        yield sse({"step": 1, "message": "Analyzing image with Claude Vision…", "status": "active"})
        await asyncio.sleep(0.05)

        log["steps"]["vision_request"] = {
            "sent_at":    _ts(),
            "model":      "claude-sonnet-4-5",
            "prompt":     LAYERED_EXTRACTION_PROMPT,
            "hint_text":  hint_text,
        }

        t_vision = time.monotonic()
        raw_text = await loop.run_in_executor(
            None,
            lambda: _call_vision_model(image_data, mime_type, LAYERED_EXTRACTION_PROMPT, hint_text, api_key),
        )
        vision_ms = round((time.monotonic() - t_vision) * 1000)

        log["steps"]["vision_response"] = {
            "received_at": _ts(),
            "duration_ms": vision_ms,
            "raw_text":    raw_text,
        }
        logger.info("job=%s Claude responded in %dms (%d chars)",
                    job_id, vision_ms, len(raw_text))

        yield sse({"step": 1, "message": "Analyzing image with Claude Vision…", "status": "done"})

        # ── Step 2: Strip markdown fencing ────────────────────────────────────
        yield sse({"step": 2, "message": "Extracting slide elements…", "status": "active"})
        await asyncio.sleep(0.05)

        if raw_text.startswith("```"):
            lines    = raw_text.split("\n")
            raw_text = "\n".join(lines[1:-1] if lines[-1].strip() == "```" else lines[1:])

        yield sse({"step": 2, "message": "Extracting slide elements…", "status": "done"})

        # ── Step 3: Parse JSON ─────────────────────────────────────────────────
        yield sse({"step": 3, "message": "Designing slide structure…", "status": "active"})
        await asyncio.sleep(0.05)

        slide_data = json.loads(raw_text)
        slide_data["layout"] = "layered"

        log["steps"]["json_parse"] = {
            "parsed_at":  _ts(),
            "slide_data": slide_data,   # full raw structure before any rendering
        }
        logger.info("job=%s parsed %d elements", job_id, len(slide_data.get("elements", [])))

        preview = _build_preview(slide_data)
        yield sse({"step": 3, "message": "Designing slide structure…", "status": "done",
                   "preview": preview})

        # ── Step 4: Build PPTX ─────────────────────────────────────────────────
        yield sse({"step": 4, "message": "Building PowerPoint file…", "status": "active"})
        await asyncio.sleep(0.05)

        t_build = time.monotonic()
        pptx_bytes, used_fallback = await loop.run_in_executor(
            None, lambda: build_pptx(slide_data, theme_name)
        )
        build_ms = round((time.monotonic() - t_build) * 1000)

        log["steps"]["pptx_build"] = {
            "built_at":     _ts(),
            "duration_ms":  build_ms,
            "used_fallback": used_fallback,
            "output_bytes": len(pptx_bytes),
        }
        logger.info("job=%s pptx built in %dms fallback=%s bytes=%d",
                    job_id, build_ms, used_fallback, len(pptx_bytes))

        token = str(uuid.uuid4())
        downloads[token] = pptx_bytes
        jobs.pop(job_id, None)

        yield sse({"step": 4, "message": "Building PowerPoint file…", "status": "done",
                   "download_token": token})

    except json.JSONDecodeError as exc:
        msg = f"Failed to parse Claude response: {exc}"
        log["error"] = msg
        logger.error("job=%s %s", job_id, msg)
        yield f"data: {json.dumps({'error': msg})}\n\n"

    except Exception as exc:
        log["error"] = str(exc)
        logger.error("job=%s unexpected error: %s", job_id, exc, exc_info=True)
        yield f"data: {json.dumps({'error': str(exc)})}\n\n"

    finally:
        log["completed_at"]      = _ts()
        log["total_duration_ms"] = round((time.monotonic() - t_start) * 1000)
        _write_job_log(job_id, log)


@app.get("/themes")
async def themes():
    """Return the list of available slide themes for the UI dropdown."""
    file_themes = set()
    if os.path.isdir(TEMPLATES_DIR):
        for fname in os.listdir(TEMPLATES_DIR):
            if fname.lower().endswith(".pptx"):
                file_themes.add(fname[:-5])  # strip .pptx

    result = []
    for name, palette in SLIDE_THEMES.items():
        result.append({
            "name":         name,
            "desc":         palette["desc"],
            "bg":           palette["bg"],
            "accent":       palette["accent"],
            "has_template": name in file_themes,
        })
    return {"themes": result}


@app.post("/generate")
async def generate(
    file: UploadFile = File(...),
    api_key: str = Form(...),
    hint_text: str = Form(""),
    theme: str = Form("Default"),
):
    """Accept image upload, API key, optional hint text, and theme; return job_id."""
    if not api_key or not api_key.strip():
        raise HTTPException(status_code=400, detail="API key is required")
    if not file.content_type or not file.content_type.startswith("image/"):
        raise HTTPException(status_code=400, detail="File must be an image")
    if theme not in SLIDE_THEMES:
        theme = "Default"

    image_data = await file.read()
    job_id = str(uuid.uuid4())
    jobs[job_id] = {
        "image_data": image_data,
        "mime_type":  file.content_type,
        "api_key":    api_key.strip(),
        "hint_text":  hint_text.strip(),
        "theme":      theme,
    }

    # Save the original image to logs/<job_id>/image.<ext>
    try:
        log_dir   = _job_log_dir(job_id)
        img_path  = os.path.join(log_dir, f"image{_mime_to_ext(file.content_type)}")
        with open(img_path, "wb") as fh:
            fh.write(image_data)
        logger.info("job=%s image saved to %s (%d bytes)", job_id, img_path, len(image_data))
    except Exception as exc:
        logger.warning("job=%s failed to save image: %s", job_id, exc)

    return {"job_id": job_id}


@app.get("/stream/{job_id}")
async def stream(job_id: str):
    """SSE stream for a job."""
    return StreamingResponse(
        process_job(job_id),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "X-Accel-Buffering": "no",
        },
    )


@app.get("/download/{token}")
async def download(token: str):
    """Download the generated .pptx file."""
    pptx_bytes = downloads.get(token)
    if not pptx_bytes:
        raise HTTPException(status_code=404, detail="File not found or already downloaded")

    del downloads[token]

    return StreamingResponse(
        BytesIO(pptx_bytes),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": "attachment; filename=slide.pptx"},
    )


app.mount("/", StaticFiles(directory="static", html=True), name="static")
